"""File service for handling uploads and document processing."""

import os
import uuid
import logging
import mimetypes
import asyncio
from typing import Dict, List, Any, Optional, Tuple, BinaryIO
from datetime import datetime
import aiofiles
try:
    import magic
except ImportError:
    magic = None

# Document processing imports
import PyPDF2
from docx import Document
from pptx import Presentation
import pandas as pd
from PIL import Image
import io

from ..utils.openai_client import openai_client
from ..utils.supabase_client import supabase_client
from ..config.settings import settings
from ..models.schemas import (
    FileUploadResponse,
    FileResponse,
    FileProcessRequest,
    FileType
)

logger = logging.getLogger(__name__)


class FileService:
    """Service for handling file uploads and document processing."""
    
    def __init__(self):
        """Initialize File Service."""
        self.upload_dir = settings.upload_dir
        self.max_file_size = settings.max_file_size
        self.supported_types = settings.supported_file_types_list
        
        # Ensure upload directory exists
        os.makedirs(self.upload_dir, exist_ok=True)
    
    async def upload_file(
        self,
        file_content: bytes,
        filename: str,
        content_type: str,
        user_id: str,
        assistant_id: str = None,
        metadata: Dict[str, Any] = None
    ) -> FileUploadResponse:
        """Upload and process a file."""
        try:
            # Validate file
            await self._validate_file(file_content, filename, content_type)
            
            # Generate unique filename
            file_id = str(uuid.uuid4())
            file_extension = os.path.splitext(filename)[1].lower()
            unique_filename = f"{file_id}{file_extension}"
            file_path = os.path.join(self.upload_dir, unique_filename)
            
            # Save file to disk
            async with aiofiles.open(file_path, 'wb') as f:
                await f.write(file_content)
            
            # Create database record
            db_file = await supabase_client.create_file_record(
                user_id=user_id,
                filename=unique_filename,
                original_filename=filename,
                file_type=content_type,
                file_size=len(file_content),
                upload_path=file_path,
                assistant_id=assistant_id,
                metadata=metadata or {}
            )
            
            logger.info(f"Uploaded file: {db_file['id']} ({filename})")
            
            # Start background processing
            asyncio.create_task(self._process_file_background(db_file))
            
            return FileUploadResponse(
                id=db_file["id"],
                filename=unique_filename,
                original_filename=filename,
                file_type=content_type,
                file_size=len(file_content),
                upload_path=file_path,
                processing_status="pending",
                created_at=datetime.fromisoformat(db_file["created_at"]),
                metadata=metadata or {}
            )
            
        except Exception as e:
            logger.error(f"Error uploading file {filename}: {str(e)}")
            raise
    
    async def process_files_for_assistant(
        self,
        request: FileProcessRequest,
        user_id: str
    ) -> Dict[str, Any]:
        """Process files for a specific assistant."""
        try:
            processed_files = []
            failed_files = []
            
            for file_id in request.file_ids:
                try:
                    # Get file record
                    files = await supabase_client.get_user_files(user_id)
                    file_record = None
                    for f in files:
                        if f["id"] == file_id:
                            file_record = f
                            break
                    
                    if not file_record:
                        failed_files.append({"file_id": file_id, "error": "File not found"})
                        continue
                    
                    # Process and upload to OpenAI
                    openai_file_id = await self._upload_to_openai(file_record)
                    
                    # Update file record
                    await supabase_client.update_file_record(
                        file_id,
                        {
                            "openai_file_id": openai_file_id,
                            "assistant_id": request.assistant_id,
                            "processing_status": "processed"
                        }
                    )
                    
                    processed_files.append({
                        "file_id": file_id,
                        "openai_file_id": openai_file_id,
                        "filename": file_record["original_filename"]
                    })
                    
                except Exception as e:
                    logger.error(f"Error processing file {file_id}: {str(e)}")
                    failed_files.append({"file_id": file_id, "error": str(e)})
            
            # Add processed files to assistant's vector store
            if processed_files:
                from ..services.course_agent import course_agent_service
                
                openai_file_ids = [f["openai_file_id"] for f in processed_files]
                await course_agent_service.add_files_to_course_assistant(
                    request.assistant_id, openai_file_ids, user_id
                )
            
            return {
                "assistant_id": request.assistant_id,
                "processed_files": processed_files,
                "failed_files": failed_files,
                "total_processed": len(processed_files),
                "total_failed": len(failed_files)
            }
            
        except Exception as e:
            logger.error(f"Error processing files for assistant: {str(e)}")
            raise
    
    async def get_user_files(
        self,
        user_id: str,
        assistant_id: str = None,
        limit: int = 50
    ) -> List[FileResponse]:
        """Get user's uploaded files."""
        try:
            files = await supabase_client.get_user_files(user_id, assistant_id, limit)
            
            result = []
            for file_record in files:
                result.append(FileResponse(
                    id=file_record["id"],
                    filename=file_record["filename"],
                    original_filename=file_record["original_filename"],
                    file_type=file_record["file_type"],
                    file_size=file_record["file_size"],
                    upload_path=file_record["upload_path"],
                    openai_file_id=file_record.get("openai_file_id"),
                    processing_status=file_record.get("processing_status", "pending"),
                    error_message=file_record.get("error_message"),
                    assistant_id=file_record.get("assistant_id"),
                    vector_store_id=file_record.get("vector_store_id"),
                    chunk_count=file_record.get("chunk_count", 0),
                    uploaded_by=file_record["uploaded_by"],
                    created_at=datetime.fromisoformat(file_record["created_at"]),
                    updated_at=datetime.fromisoformat(file_record["updated_at"]) if file_record.get("updated_at") else None,
                    metadata=file_record.get("metadata", {})
                ))
            
            return result
            
        except Exception as e:
            logger.error(f"Error getting user files: {str(e)}")
            raise
    
    async def delete_file(
        self,
        file_id: str,
        user_id: str
    ) -> bool:
        """Delete a file."""
        try:
            # Get file record
            files = await supabase_client.get_user_files(user_id)
            file_record = None
            for f in files:
                if f["id"] == file_id:
                    file_record = f
                    break
            
            if not file_record:
                raise Exception("File not found")
            
            # Delete from OpenAI if uploaded
            if file_record.get("openai_file_id"):
                try:
                    await openai_client.delete_file(file_record["openai_file_id"])
                except Exception as e:
                    logger.warning(f"Failed to delete OpenAI file: {str(e)}")
            
            # Delete local file
            try:
                if os.path.exists(file_record["upload_path"]):
                    os.remove(file_record["upload_path"])
            except Exception as e:
                logger.warning(f"Failed to delete local file: {str(e)}")
            
            # Delete database record
            await supabase_client.delete_file_record(file_id, user_id)
            
            logger.info(f"Deleted file: {file_id}")
            return True
            
        except Exception as e:
            logger.error(f"Error deleting file {file_id}: {str(e)}")
            raise
    
    async def extract_text_content(
        self,
        file_path: str,
        file_type: str,
        chunk_size: int = 1000,
        chunk_overlap: int = 200
    ) -> List[Dict[str, Any]]:
        """Extract text content from file and create chunks."""
        try:
            text_content = ""
            
            if file_type == FileType.PDF:
                text_content = await self._extract_pdf_text(file_path)
            elif file_type == FileType.DOCX:
                text_content = await self._extract_docx_text(file_path)
            elif file_type == FileType.PPTX:
                text_content = await self._extract_pptx_text(file_path)
            elif file_type == FileType.CSV:
                text_content = await self._extract_csv_text(file_path)
            elif file_type == FileType.TXT:
                text_content = await self._extract_txt_text(file_path)
            elif file_type in [FileType.JPEG, FileType.PNG, FileType.GIF]:
                # For images, we could use OCR here
                text_content = f"Image file: {os.path.basename(file_path)}"
            else:
                raise Exception(f"Unsupported file type: {file_type}")
            
            # Create text chunks
            chunks = self._create_text_chunks(text_content, chunk_size, chunk_overlap)
            
            return chunks
            
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {str(e)}")
            raise
    
    async def _validate_file(
        self,
        file_content: bytes,
        filename: str,
        content_type: str
    ) -> None:
        """Validate uploaded file."""
        # Check file size
        if len(file_content) > self.max_file_size:
            raise Exception(f"File too large. Maximum size: {self.max_file_size / (1024*1024):.1f}MB")
        
        # Check file type
        if content_type not in self.supported_types:
            raise Exception(f"Unsupported file type: {content_type}")
        
        # Verify content type matches file content
        if magic is not None:
            try:
                detected_type = magic.from_buffer(file_content, mime=True)
                if detected_type not in self.supported_types:
                    logger.warning(f"Detected type {detected_type} differs from declared {content_type}")
            except Exception as e:
                logger.warning(f"Could not detect file type: {str(e)}")
        else:
            logger.warning("python-magic not available, skipping file type detection")
        
        # Check filename
        if not filename or len(filename) > 255:
            raise Exception("Invalid filename")
        
        # Check for potentially dangerous file extensions
        dangerous_extensions = ['.exe', '.bat', '.cmd', '.scr', '.com', '.pif']
        file_extension = os.path.splitext(filename)[1].lower()
        if file_extension in dangerous_extensions:
            raise Exception("Dangerous file type not allowed")
    
    async def _process_file_background(self, file_record: Dict[str, Any]) -> None:
        """Process file in background after upload."""
        try:
            file_id = file_record["id"]
            
            # Update status to processing
            await supabase_client.update_file_record(
                file_id,
                {"processing_status": "processing"}
            )
            
            # Extract text content (for future vector search)
            try:
                chunks = await self.extract_text_content(
                    file_record["upload_path"],
                    file_record["file_type"]
                )
                
                # Update with chunk count
                await supabase_client.update_file_record(
                    file_id,
                    {
                        "processing_status": "completed",
                        "chunk_count": len(chunks)
                    }
                )
                
                logger.info(f"Processed file {file_id}: {len(chunks)} chunks")
                
            except Exception as e:
                error_msg = f"Text extraction failed: {str(e)}"
                await supabase_client.update_file_record(
                    file_id,
                    {
                        "processing_status": "failed",
                        "error_message": error_msg
                    }
                )
                logger.error(f"Failed to process file {file_id}: {error_msg}")
                
        except Exception as e:
            logger.error(f"Background processing error for file {file_record['id']}: {str(e)}")
    
    async def _upload_to_openai(self, file_record: Dict[str, Any]) -> str:
        """Upload file to OpenAI."""
        try:
            file_path = file_record["upload_path"]
            
            # Upload to OpenAI
            openai_file = await openai_client.upload_file(
                file_path=file_path,
                purpose="assistants"
            )
            
            logger.info(f"Uploaded file to OpenAI: {openai_file.id}")
            return openai_file.id
            
        except Exception as e:
            logger.error(f"Error uploading to OpenAI: {str(e)}")
            raise
    
    async def _extract_pdf_text(self, file_path: str) -> str:
        """Extract text from PDF file."""
        try:
            text_content = ""
            
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                for page in pdf_reader.pages:
                    text_content += page.extract_text() + "\n"
            
            return text_content.strip()
            
        except Exception as e:
            logger.error(f"Error extracting PDF text: {str(e)}")
            raise
    
    async def _extract_docx_text(self, file_path: str) -> str:
        """Extract text from DOCX file."""
        try:
            doc = Document(file_path)
            text_content = ""
            
            for paragraph in doc.paragraphs:
                text_content += paragraph.text + "\n"
            
            return text_content.strip()
            
        except Exception as e:
            logger.error(f"Error extracting DOCX text: {str(e)}")
            raise
    
    async def _extract_pptx_text(self, file_path: str) -> str:
        """Extract text from PPTX file."""
        try:
            prs = Presentation(file_path)
            text_content = ""
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"
            
            return text_content.strip()
            
        except Exception as e:
            logger.error(f"Error extracting PPTX text: {str(e)}")
            raise
    
    async def _extract_csv_text(self, file_path: str) -> str:
        """Extract text from CSV file."""
        try:
            df = pd.read_csv(file_path)
            
            # Convert DataFrame to text representation
            text_content = f"CSV Data ({len(df)} rows, {len(df.columns)} columns)\n\n"
            text_content += f"Columns: {', '.join(df.columns)}\n\n"
            text_content += df.to_string(index=False)
            
            return text_content
            
        except Exception as e:
            logger.error(f"Error extracting CSV text: {str(e)}")
            raise
    
    async def _extract_txt_text(self, file_path: str) -> str:
        """Extract text from TXT file."""
        try:
            async with aiofiles.open(file_path, 'r', encoding='utf-8') as file:
                text_content = await file.read()
            
            return text_content
            
        except Exception as e:
            logger.error(f"Error extracting TXT text: {str(e)}")
            raise
    
    def _create_text_chunks(
        self,
        text: str,
        chunk_size: int = 1000,
        chunk_overlap: int = 200
    ) -> List[Dict[str, Any]]:
        """Create overlapping text chunks."""
        try:
            chunks = []
            text_length = len(text)
            
            if text_length <= chunk_size:
                chunks.append({
                    "chunk_id": str(uuid.uuid4()),
                    "content": text,
                    "start_index": 0,
                    "end_index": text_length,
                    "chunk_number": 1
                })
                return chunks
            
            start = 0
            chunk_number = 1
            
            while start < text_length:
                end = min(start + chunk_size, text_length)
                
                # Try to break at word boundaries
                if end < text_length:
                    # Look for the last space within the chunk
                    last_space = text.rfind(' ', start, end)
                    if last_space > start:
                        end = last_space
                
                chunk_content = text[start:end].strip()
                
                if chunk_content:
                    chunks.append({
                        "chunk_id": str(uuid.uuid4()),
                        "content": chunk_content,
                        "start_index": start,
                        "end_index": end,
                        "chunk_number": chunk_number
                    })
                    chunk_number += 1
                
                # Move start position with overlap
                start = max(start + chunk_size - chunk_overlap, end)
            
            return chunks
            
        except Exception as e:
            logger.error(f"Error creating text chunks: {str(e)}")
            raise
    
    async def get_file_statistics(self, user_id: str) -> Dict[str, Any]:
        """Get file upload statistics for a user."""
        try:
            files = await supabase_client.get_user_files(user_id, limit=1000)
            
            total_files = len(files)
            total_size = sum(f.get("file_size", 0) for f in files)
            
            # Count by file type
            type_counts = {}
            for file_record in files:
                file_type = file_record.get("file_type", "unknown")
                type_counts[file_type] = type_counts.get(file_type, 0) + 1
            
            # Count by processing status
            status_counts = {}
            for file_record in files:
                status = file_record.get("processing_status", "unknown")
                status_counts[status] = status_counts.get(status, 0) + 1
            
            return {
                "user_id": user_id,
                "total_files": total_files,
                "total_size_mb": round(total_size / (1024 * 1024), 2),
                "file_types": type_counts,
                "processing_status": status_counts,
                "recent_uploads": len([
                    f for f in files 
                    if (datetime.utcnow() - datetime.fromisoformat(f["created_at"])).days <= 7
                ])
            }
            
        except Exception as e:
            logger.error(f"Error getting file statistics: {str(e)}")
            raise


# Global instance
file_service = FileService()
